#!/usr/bin/env python3
"""
Update PowerPoint poster with real application screenshots
"""
from pptx import Presentation
from pptx.util import Inches
import os
import shutil

# File paths
PPTX_FILE = "StockFlow_Pitch_Deck.pptx"
BACKUP_FILE = "StockFlow_Pitch_Deck_backup.pptx"
SCREENSHOTS_DIR = "docs/screenshots"

# Screenshot mapping
screenshots = {
    "01_dashboard": "01_dashboard_1765989751219.png",
    "02_stock_level": "02_stock_level_1765989794917.png",
    "03_inventory": "03_inventory_table_1765989919002.png",
    "04_analytics": "04_analytics_1765989960070.png",
    "05_settings": "05_settings_1765990030067.png",
    "06_alerts": "06_alerts_1765990092950.png",
}

def update_presentation():
    """Update the PowerPoint presentation with real screenshots"""
    
    # Create backup
    if os.path.exists(PPTX_FILE):
        shutil.copy2(PPTX_FILE, BACKUP_FILE)
        print(f"✅ Created backup: {BACKUP_FILE}")
    
    # Load presentation
    prs = Presentation(PPTX_FILE)
    print(f"📄 Loaded presentation with {len(prs.slides)} slides")
    
    # Track updates
    updates_made = 0
    
    # Iterate through slides and add/replace screenshots
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        print(f"\n📍 Processing Slide {slide_num}...")
        
        # Look for existing images or placeholders
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                print(f"  Found image: {shape.name}")
            if hasattr(shape, "text") and shape.text:
                print(f"  Text: {shape.text[:50]}...")
        
        # Add screenshots to specific slides based on content
        # Slide 3-4: Product/Feature slides - add dashboard and inventory screenshots
        if slide_num == 3:
            # Add dashboard screenshot
            img_path = os.path.join(SCREENSHOTS_DIR, screenshots["01_dashboard"])
            if os.path.exists(img_path):
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                slide.shapes.add_picture(img_path, left, top, width=width)
                print(f"  ✅ Added dashboard screenshot")
                updates_made += 1
        
        elif slide_num == 4:
            # Add inventory screenshot
            img_path = os.path.join(SCREENSHOTS_DIR, screenshots["03_inventory"])
            if os.path.exists(img_path):
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                slide.shapes.add_picture(img_path, left, top, width=width)
                print(f"  ✅ Added inventory screenshot")
                updates_made += 1
        
        elif slide_num == 5:
            # Add analytics screenshot
            img_path = os.path.join(SCREENSHOTS_DIR, screenshots["04_analytics"])
            if os.path.exists(img_path):
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                slide.shapes.add_picture(img_path, left, top, width=width)
                print(f"  ✅ Added analytics screenshot")
                updates_made += 1
    
    # Save the updated presentation
    output_file = "StockFlow_Pitch_Deck_Updated.pptx"
    prs.save(output_file)
    print(f"\n✅ Saved updated presentation: {output_file}")
    print(f"📊 Total updates made: {updates_made}")
    
    # Also create a version with all screenshots in one slide for reference
    create_screenshot_showcase()

def create_screenshot_showcase():
    """Create a new slide with all screenshots showcased"""
    prs = Presentation(PPTX_FILE)
    
    # Add a new slide at the end
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "Digital Inventory - Application Screenshots"
    
    # Add screenshots in a grid (2x3)
    screenshots_list = [
        ("01_dashboard", "Dashboard"),
        ("02_stock_level", "Stock Level"),
        ("03_inventory", "Inventory"),
        ("04_analytics", "Analytics"),
        ("05_settings", "Settings"),
        ("06_alerts", "Alerts"),
    ]
    
    row, col = 0, 0
    for screenshot_key, label in screenshots_list:
        img_path = os.path.join(SCREENSHOTS_DIR, screenshots[screenshot_key])
        
        if os.path.exists(img_path):
            # Calculate position
            left = Inches(0.5 + col * 3.5)
            top = Inches(1.2 + row * 2.5)
            width = Inches(3)
            
            # Add image
            slide.shapes.add_picture(img_path, left, top, width=width)
            
            # Add label
            label_box = slide.shapes.add_textbox(left, top + Inches(1.8), width, Inches(0.3))
            label_frame = label_box.text_frame
            label_frame.text = label
            
            # Move to next position
            col += 1
            if col >= 2:
                col = 0
                row += 1
    
    # Save with screenshots showcase
    output_file = "StockFlow_with_Screenshots.pptx"
    prs.save(output_file)
    print(f"✅ Created screenshot showcase: {output_file}")

if __name__ == "__main__":
    print("🚀 Starting PowerPoint Update...\n")
    update_presentation()
    print("\n✨ Done!")
